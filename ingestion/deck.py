import fitz
from pptx import Presentation
from pathlib import Path


def extract_pdf(path):
    doc = fitz.open(path)
    chunks = []
    for i, page in enumerate(doc):
        text = page.get_text().strip()
        if text:
            chunks.append({
                "text": text,
                "source": "deck",
                "reference": f"slide_{i + 1}"
            })
    return chunks


def extract_pptx(path):
    prs = Presentation(path)
    chunks = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            chunks.append({
                "text": "\n".join(texts),
                "source": "deck",
                "reference": f"slide_{i + 1}"
            })
    return chunks


def extract_txt(path):
    text = Path(path).read_text(encoding="utf-8").strip()
    if not text:
        return []
    sections = [part.strip() for part in text.split("\n\n") if part.strip()]
    return [
        {
            "text": section,
            "source": "deck",
            "reference": f"slide_{i + 1}"
        }
        for i, section in enumerate(sections)
    ]


def ingest(path):
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        return extract_pdf(path)
    elif ext in (".pptx", ".ppt"):
        return extract_pptx(path)
    elif ext == ".txt":
        return extract_txt(path)
    return []
